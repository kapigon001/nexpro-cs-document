"""Builder Agent - Responsible for actual PowerPoint file generation"""
from typing import Any, Dict, List, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from .base_agent import BaseAgent
from ..core.task import Task


class BuilderAgent(BaseAgent):
    """
    Builder Agent: Generates the actual PowerPoint file.

    Responsibilities:
    - Create PowerPoint presentation files
    - Apply designs and layouts
    - Insert content into slides
    - Handle charts and images
    - Save output files
    """

    def __init__(self, output_dir: str = "output"):
        super().__init__(
            name="BuilderAgent",
            role="PowerPoint Generation",
            description="Generates PowerPoint files using python-pptx library"
        )
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    async def execute_task(self, task: Task) -> Any:
        """Execute build-related tasks"""
        task_type = task.input_data.get("type", "build_presentation") if task.input_data else "build_presentation"

        if task_type == "build_presentation":
            return await self._build_presentation(task.input_data)
        elif task_type == "build_slide":
            return await self._build_slide(task.input_data)
        elif task_type == "save_presentation":
            return await self._save_presentation(task.input_data)
        else:
            return await self._build_presentation(task.input_data)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Convert hex color to RGBColor"""
        hex_color = hex_color.lstrip('#')
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    async def _build_presentation(self, input_data: Dict) -> Dict:
        """Build complete PowerPoint presentation"""
        content = input_data.get("content", {})
        design = input_data.get("design", {})
        filename = input_data.get("filename", "presentation.pptx")

        self.log("Building PowerPoint presentation...")

        # Create presentation
        prs = Presentation()

        # Set slide dimensions (16:9)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        # Get theme colors
        theme = design.get("theme", {})
        colors = theme.get("colors", {
            "primary": "1F4E79",
            "text": "333333",
            "background": "FFFFFF",
        })

        # Build each slide
        for i, slide_content in enumerate(content.get("slides", [])):
            slide_design = None
            if "slides" in design and i < len(design["slides"]):
                slide_design = design["slides"][i]

            await self._add_slide(prs, slide_content, slide_design, colors, theme)

        # Save presentation
        output_path = self.output_dir / filename
        prs.save(str(output_path))

        self.log(f"Presentation saved to: {output_path}")

        return {
            "success": True,
            "file_path": str(output_path),
            "slide_count": len(content.get("slides", [])),
        }

    async def _add_slide(
        self,
        prs: Presentation,
        content: Dict,
        design: Optional[Dict],
        colors: Dict,
        theme: Dict
    ):
        """Add a single slide to the presentation"""
        slide_type = content.get("type", "content")

        # Use blank layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Get font settings
        fonts = theme.get("fonts", {})
        title_font = fonts.get("title", {"name": "Yu Gothic UI", "size": 36, "bold": True})
        body_font = fonts.get("body", {"name": "Yu Gothic UI", "size": 18, "bold": False})

        if slide_type == "title":
            await self._build_title_slide(slide, content, colors, title_font)
        elif slide_type == "agenda":
            await self._build_content_slide(slide, content, colors, title_font, body_font)
        elif slide_type in ["content", "conclusion"]:
            await self._build_content_slide(slide, content, colors, title_font, body_font)
        elif slide_type in ["two_column", "comparison"]:
            await self._build_two_column_slide(slide, content, colors, title_font, body_font)
        else:
            await self._build_content_slide(slide, content, colors, title_font, body_font)

    async def _build_title_slide(self, slide, content: Dict, colors: Dict, title_font: Dict):
        """Build a title slide"""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.get("title", "")
        title_para.font.size = Pt(title_font.get("size", 44))
        title_para.font.bold = title_font.get("bold", True)
        title_para.font.color.rgb = self._hex_to_rgb(colors.get("primary", "1F4E79"))
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        if content.get("subtitle"):
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.5), Inches(9), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = content.get("subtitle", "")
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = self._hex_to_rgb(colors.get("text", "333333"))
            subtitle_para.alignment = PP_ALIGN.CENTER

    async def _build_content_slide(
        self, slide, content: Dict, colors: Dict, title_font: Dict, body_font: Dict
    ):
        """Build a content slide with bullet points"""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.get("title", "")
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb(colors.get("primary", "1F4E79"))

        # Body content
        body_items = content.get("body", content.get("items", []))
        if body_items:
            body_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(9), Inches(4)
            )
            body_frame = body_box.text_frame
            body_frame.word_wrap = True

            for i, item in enumerate(body_items):
                if i == 0:
                    para = body_frame.paragraphs[0]
                else:
                    para = body_frame.add_paragraph()

                para.text = f"• {item}"
                para.font.size = Pt(body_font.get("size", 18))
                para.font.color.rgb = self._hex_to_rgb(colors.get("text", "333333"))
                para.space_after = Pt(12)

    async def _build_two_column_slide(
        self, slide, content: Dict, colors: Dict, title_font: Dict, body_font: Dict
    ):
        """Build a two-column slide"""
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = content.get("title", "")
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb(colors.get("primary", "1F4E79"))

        # Left column
        left_items = content.get("left", [])
        if left_items:
            left_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(4.2), Inches(4)
            )
            left_frame = left_box.text_frame
            left_frame.word_wrap = True

            for i, item in enumerate(left_items):
                if i == 0:
                    para = left_frame.paragraphs[0]
                else:
                    para = left_frame.add_paragraph()
                para.text = f"• {item}"
                para.font.size = Pt(16)
                para.font.color.rgb = self._hex_to_rgb(colors.get("text", "333333"))

        # Right column
        right_items = content.get("right", [])
        if right_items:
            right_box = slide.shapes.add_textbox(
                Inches(5.2), Inches(1.3), Inches(4.2), Inches(4)
            )
            right_frame = right_box.text_frame
            right_frame.word_wrap = True

            for i, item in enumerate(right_items):
                if i == 0:
                    para = right_frame.paragraphs[0]
                else:
                    para = right_frame.add_paragraph()
                para.text = f"• {item}"
                para.font.size = Pt(16)
                para.font.color.rgb = self._hex_to_rgb(colors.get("text", "333333"))

    async def _build_slide(self, input_data: Dict) -> Dict:
        """Build a single slide (for incremental building)"""
        # This would be used for building slides one at a time
        self.log("Building individual slide...")
        return {"success": True, "message": "Slide built"}

    async def _save_presentation(self, input_data: Dict) -> Dict:
        """Save presentation to file"""
        presentation = input_data.get("presentation")
        filename = input_data.get("filename", "output.pptx")

        if presentation is None:
            raise ValueError("No presentation object provided")

        output_path = self.output_dir / filename
        presentation.save(str(output_path))

        self.log(f"Saved presentation: {output_path}")
        return {"success": True, "file_path": str(output_path)}
